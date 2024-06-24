import pyperclip
from pptx import Presentation
from pptx.util import Pt
import os
import string
from tkinter import Tk, filedialog

def open_presentation(file_path):
    return Presentation(file_path)

def create_presentation():
    return Presentation()

def add_slide_with_text(presentation, text):
    slide_layout = presentation.slide_layouts[1]  # Use layout 1 for content
    slide = presentation.slides.add_slide(slide_layout)
    subtitle = slide.placeholders[1]
    clean_text = text.replace('_x00D_', '').replace('\r\n', '\n')
    subtitle.text = clean_text
    
    # Set font size and name for all runs in the subtitle
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.name = "Times New Roman"

def save_presentation(presentation, file_name):
    presentation.save(file_name)

def validate_file_name(file_name):
    # Validate the file name to ensure it's not empty and doesn't contain illegal characters
    valid_chars = f"-_.() {string.ascii_letters}{string.digits}"
    return all(char in valid_chars for char in file_name)

def ask_user_for_option():
    print("Welcome to PowerPoint Presentation Manager!")
    print("1. Create a new presentation")
    print("2. Append to an existing presentation")
    while True:
        choice = input("Please enter your choice (1 or 2): ")
        if choice in ['1', '2']:
            return choice
        else:
            print("Invalid choice. Please enter 1 or 2.")

def choose_file_dialog():
    root = Tk()
    root.withdraw()  # Hide the main window
    file_path = filedialog.askopenfilename(title="Choose PowerPoint file", filetypes=[("PowerPoint files", "*.pptx")])
    return file_path

def main():
    try:
        choice = ask_user_for_option()
        if choice == '1':
            presentation = create_presentation()
            print("New presentation created.")
            file_name = input("Enter the file name to save the presentation (without extension): ")
            while not validate_file_name(file_name):
                print("Invalid file name. Please provide a valid file name.")
                file_name = input("Enter the file name to save the presentation (without extension): ")
            file_path = f"{file_name}.pptx"
        elif choice == '2':
            file_path = choose_file_dialog()
            while not os.path.exists(file_path):
                print("File not found. Please choose a valid file.")
                file_path = choose_file_dialog()
            presentation = open_presentation(file_path)
            print("Existing presentation opened.")
        
        while True:
            print("Waiting for new text from clipboard...")
            new_text = pyperclip.waitForNewPaste()
            print("Text copied from clipboard.")
            add_slide_with_text(presentation, new_text)
            print("Text added to new slide!")
    except KeyboardInterrupt:
        print("\nKeyboard interrupt detected. Saving presentation...")
        save_presentation(presentation, file_path)
        print(f"Presentation saved as {file_path}")
    except Exception as e:
        print(f"An error occurred: {e}")

if __name__ == "__main__":
    main()
